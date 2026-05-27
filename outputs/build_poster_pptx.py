"""Build an EDITABLE 48x36in PowerPoint of the 'Edge or Beta?' poster.

Text is native editable text boxes; charts are the existing PNGs. Layout maps the
HTML grid at 96 px/in. Fonts (Fraunces / IBM Plex) fall back to similar faces if
not installed. Run: python outputs/build_poster_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
from pathlib import Path

HERE = Path(__file__).resolve().parent
FIG = HERE.parent / "doc" / "202605" / "25"
OUT = FIG / "poster_editable.pptx"

# palette
NAVY="0B2545"; INK="0E1726"; MUTED="5A6678"; HAIR="D7D2C6"; PAPER="F4F1EA"
CARD="FFFFFF"; CYAN="0E7490"; AMBER="B45309"; RED="A41E22"; GREEN="15803D"; GOLD="C7A100"
ICE="EDF2F9"; ICEB="CFE0F0"; TINT="FBF3E8"; GREENT="EAF4EC"
SERIF="Fraunces"; SANS="IBM Plex Sans"; MONO="IBM Plex Mono"
def C(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width = Inches(48); prs.slide_height = Inches(36)
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def rect(x,y,w,h,fill=None,line=None,lw=1.0,round=False,shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
                             Inches(x),Inches(y),Inches(w),Inches(h))
    if round:
        try: shp.adjustments[0]=0.04
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=C(fill)
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=C(line); shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp

def text(x,y,w,h,paras,anchor=MSO_ANCHOR.TOP,align=PP_ALIGN.LEFT,wrap=True):
    """paras = list of paragraphs; each para = list of runs (text,size,color,bold,italic,font)."""
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        if isinstance(para,dict):
            runs=para["runs"];
            if "space_after" in para: p.space_after=Pt(para["space_after"])
            if "line" in para:
                from pptx.oxml.ns import qn as _q
        else:
            runs=para
        for (t,sz,col,*rest) in runs:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.color.rgb=C(col)
            r.font.bold=rest[0] if len(rest)>0 else False
            r.font.italic=rest[1] if len(rest)>1 else False
            r.font.name=rest[2] if len(rest)>2 else SANS
    return tb

def pic(path,x,y,bw,bh=None):
    """Place image fit within box bw x bh (contain, centered). bh=None -> width-bound."""
    im=Image.open(FIG/path); ar=im.height/im.width
    if bh is None:
        w=bw; h=bw*ar; px=x; py=y
    elif bw*ar<=bh:
        w=bw; h=bw*ar; px=x; py=y+(bh-h)/2
    else:
        h=bh; w=bh/ar; px=x+(bw-w)/2; py=y
    s.shapes.add_picture(str(FIG/path),Inches(px),Inches(py),Inches(w),Inches(h))
    return h

# ---------- background ----------
rect(0,0,48,36,fill=PAPER)

M=0.79; CW=10.99; DW=23.6; EW=10.99
CX=M; DX=12.13; EX=36.08

# ================= A. TITLE =================
text(M,0.34,30,0.5,[[("STOCK-MARKET RESEARCH · S&P 500 EQUITIES · 2015–2026",17,AMBER,True,False,MONO)]])
text(M,0.62,30,2.0,[[("Edge or ",95,NAVY,True,False,SERIF),("Beta?",95,RED,True,True,SERIF)]])
text(M,2.35,33,1.1,[[("In plain terms: can a clever rule for picking oversold S&P 500 stocks beat just owning the market — real skill (an ",24,INK,False,True,SERIF),
                      ("edge",24,INK,True,True,SERIF),(") or only the market's rise (",24,INK,False,True,SERIF),("beta",24,INK,True,True,SERIF),
                      (")? We tested one honestly across 11 years: ",24,INK,False,True,SERIF),("it can't.",24,RED,True,True,SERIF)]])
# byline (right)
text(34.5,0.5,12.7,2.8,[
    [("Diego Bicieg Vázquez del Mercado · Ícaro Dos Santos Rabelo",14,INK,True)],
    [("Sai Bhaskar Chundru · Osmond Xin (Yi)",14,INK,True)],
    [("DAMO-699 · M.S. Data Analytics · University of Niagara Falls Canada",13,MUTED)],
    [("Capstone Project · Spring 2026 · Supervisor: William Pourmajidi",13,MUTED)],
],align=PP_ALIGN.RIGHT)
try:
    s.shapes.add_picture(str(FIG/"unfc-logo.svg"),Inches(43.5),Inches(0.5),Inches(3.6))
except Exception: pass
rect(M,3.35,46.4,0.04,fill=NAVY)  # navy rule

# ================= B. SPINE =================
spine=[("1","MECHANISM",NAVY,"The pressure is real","Around monthly options expiry, dealer hedging nudges stock prices — a real, documented market quirk."),
       ("2","PROXY SIGNAL",AMBER,"We made it testable","RSI<22 & ≥3 red candles → buy T+1, hold 6 days, top-3."),
       ("3","STRICT TESTS",CYAN,"Then we broke it","Costs · out-of-sample · walk-forward · ML · full-universe · random-entry MC."),
       ("4","RESULT",RED,"All three: not supported","✗H1 · ✗H2 p=0.47 · ✗H3 RSI 1/3 → beta, not alpha."),
       ("5","PUBLIC TOOL",GREEN,"The null becomes useful","“Edge or Beta?” — test any rule vs random picks, random ETF timing & holding SPY.")]
by=3.7; bh=2.95; bw=(46.4-4*0.27)/5
text(M,by-0.32,40,0.3,[[("THE RESEARCH JOURNEY  ·  mechanism → proxy → strict tests → null result → public tool",15,MUTED,False,False,MONO)]])
for i,(num,lab,col,claim,dtl) in enumerate(spine):
    bx=M+i*(bw+0.27)
    rect(bx,by,bw,bh,fill=CARD,line=HAIR,lw=1.2,round=True)
    rect(bx,by,bw,0.62,fill=col,round=True)
    text(bx+0.18,by+0.06,bw-0.3,0.5,[[(num+"   "+lab,19,"FFFFFF",True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(bx+0.18,by+0.78,bw-0.36,bh-0.9,[
        [(claim,26,INK,True,False,SERIF)],
        {"runs":[(dtl,16,MUTED)]},
    ])

# ================= GLANCE =================
glance=[("~510","S&P 500 tickers"),("2015–26","11 years tested"),("136","monthly expiries"),
        ("801","triggered signals"),("177","portfolio trades"),("~66%","base win-rate")]
gy=6.95; gh=0.95; gw=(46.4-5*0.19)/6
for i,(v,k) in enumerate(glance):
    gx=M+i*(gw+0.19)
    rect(gx,gy,gw,gh,fill=CARD,line=HAIR,lw=1.2,round=True)
    text(gx+0.15,gy+0.05,gw-0.3,gh-0.1,[[(v,30,NAVY,True,False,SERIF)],[(k,13,MUTED,False,False,MONO)]],anchor=MSO_ANCHOR.MIDDLE)

# ================= column headers =================
midy=8.2
def colhead(x,w,n,label):
    rect(x,midy,w,0.62,fill=NAVY,round=True)
    text(x+0.2,midy+0.02,w-0.3,0.58,[[(n+"   ",30,"7FB2E6",True,False,SERIF),(label,21,"FFFFFF",True,False,MONO)]],anchor=MSO_ANCHOR.MIDDLE)
colhead(CX,CW,"C","THE SETUP")
colhead(DX,DW,"D","THE TEARDOWN · THE ANALYSIS")
colhead(EX,EW,"E","SO WHAT · THE CONCLUSION")

def panel(x,y,w,h,fill=CARD,line=HAIR):
    rect(x,y,w,h,fill=fill,line=line,lw=1.2,round=True); return y

def kick(x,y,w,t,col=AMBER):
    text(x,y,w,0.35,[[(t.upper(),16,col,True,False,MONO)]]); return y+0.4

# ================= C. SETUP =================
cy=midy+0.85
# background
panel(CX,cy,CW,2.3)
yy=kick(CX+0.25,cy+0.18,CW-0.5,"Background & motivation")
text(CX+0.25,yy,CW-0.5,1.7,[[("On the third Friday each month, options expire. Dealers who sold options hedge their books; as expiry nears, charm and other hedging flows can push prices around index and option markets.",17,INK)]])
cy+=2.7
# research question + hypotheses
panel(CX,cy,CW,5.7)
yy=kick(CX+0.25,cy+0.18,CW-0.5,"Research Question & Pre-registered Hypotheses")
text(CX+0.25,yy,CW-0.5,1.2,[[("Can a simple stock rule capture an index-level mechanism — after costs?",30,NAVY,True,False,SERIF)]])
text(CX+0.25,yy+1.3,CW-0.5,0.8,[[("Is there genuine stock-selection alpha, or only market beta? Three hypotheses, fixed before testing (α = 0.05):",17,INK)]])
hyps=[("H₁","the signal beats a random-entry null"),("H₂","the edge is stronger in high VIX than low VIX"),("H₃","ML confirms RSI / red-candle features drive it")]
hyy=cy+3.1
for hl,ht in hyps:
    rect(CX+0.25,hyy,CW-0.5,0.68,fill=ICE,line=ICEB,lw=1.2,round=True)
    rect(CX+0.25,hyy,0.08,0.68,fill=NAVY)
    rect(CX+0.46,hyy+0.14,0.62,0.4,fill=NAVY,round=True)
    text(CX+0.46,hyy+0.11,0.62,0.46,[[(hl,18,"FFFFFF",True,False,SERIF)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    text(CX+1.28,hyy+0.05,CW-2.5,0.58,[[(ht,16,NAVY,True,True,SERIF)]],anchor=MSO_ANCHOR.MIDDLE)
    rect(CX+CW-0.86,hyy+0.15,0.5,0.4,fill=RED,round=True)
    text(CX+CW-0.86,hyy+0.12,0.5,0.46,[[("✗",16,"FFFFFF",True)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    hyy+=0.82
cy+=6.1
# rule + schematic
panel(CX,cy,CW,4.4)
yy=kick(CX+0.25,cy+0.18,CW-0.5,"The pre-registered rule")
text(CX+0.25,yy,CW-0.5,0.7,[[("RSI < 22  &  ≥ 3 red candles → buy T+1 open · hold 6d · top-3",18,NAVY,True,False,MONO)]])
pic("c_fig1_rule_schematic.png",CX+0.5,cy+1.4,CW-1.0,2.7)
cy+=4.8
# concepts
panel(CX,cy,CW,2.5)
yy=kick(CX+0.25,cy+0.18,CW-0.5,"Key concepts")
concepts="RSI momentum 0–100, <22 oversold  ·  VIX expected volatility (“fear index”)  ·  Expiry/charm settlement → dealer hedging flows  ·  Alpha vs Beta skill-excess vs free market return"
text(CX+0.25,yy,CW-0.5,1.8,[[(concepts,16,INK)]])
cy+=2.9
# baseline + temptation + donut
panel(CX,cy,CW,5.9)
yy=kick(CX+0.25,cy+0.18,CW-0.5,"The baseline forgotten · …and the backtest looked great",GREEN)
text(CX+0.25,yy,3.0,1.4,[[("≈12%/yr",40,AMBER,True,False,SERIF)],[("S&P 500 rose ~150% — any long-only rule inherits this for free.",14,INK)]])
tiles=[("1.35","profit factor"),("0.81","OOS composite"),("8/8","walk-fwd +")]
for j,(tv,tk) in enumerate(tiles):
    tx=CX+0.25+j*1.78
    rect(tx,yy+1.6,1.62,1.15,fill=GREENT,line="BFE0C8",lw=1.2,round=True)
    text(tx,yy+1.62,1.62,1.1,[[(tv,26,GREEN,True,False,SERIF)],[(tk,11,MUTED,False,False,MONO)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
pic("c_composite_donut.png",CX+5.6,cy+0.95,CW-5.9,4.6)
text(CX+0.25,yy+3.05,CW-0.5,1.6,[[("This is the temptation — it looks great until you remember the 12%/yr it stands on.",15,INK,True)]])

# ================= D. TEARDOWN =================
dy=midy+0.85
# Gate 1 hero (left text + right chart)
panel(DX,dy,DW,8.0)
yy=kick(DX+0.3,dy+0.2,6.0,"Gate 1 · beat the index")
text(DX+0.3,yy,6.0,1.5,[[("It never beats ",36,NAVY,True,False,SERIF),("doing nothing.",36,RED,True,False,SERIF)]])
text(DX+0.3,yy+1.7,6.0,2.5,[[("The rule (navy) stays below buy-and-hold SPY (amber) the whole period — its “profit” is market beta you could capture for free.",16,MUTED)]])
pic("fig_equity_vs_spy.png",DX+7.0,dy+0.4,DW-7.3,7.4)
dy+=8.2
# ladder
ladder=[("✗","Beat just holding the index (SPY)?","trails SPY"),
        ("✗","Beat random stock-picking?","p = 0.378"),
        ("✗","A real edge in the high-VIX regime?","p = 0.474"),
        ("✗","Do ML models rely on the rule’s features?","RSI 1/3"),
        ("=","No demonstrated alpha. The return is market beta.","β, not α")]
lh=0.62
for k,(g,q,r) in enumerate(ladder):
    ly=dy+k*(lh+0.06); fin=(g=="=")
    rect(DX,ly,DW,lh,fill=(NAVY if fin else CARD),line=HAIR,lw=1.0,round=False)
    rect(DX,ly,0.1,lh,fill=NAVY if fin else RED)
    text(DX+0.25,ly,0.7,lh,[[(g,24,("7DE0A6" if fin else RED),True,False,SERIF)]],anchor=MSO_ANCHOR.MIDDLE)
    text(DX+1.0,ly,DW-5.0,lh,[[(q,19,("FFFFFF" if fin else INK),True)]],anchor=MSO_ANCHOR.MIDDLE)
    text(DX+DW-4.0,ly,3.8,lh,[[(r,18,("7DE0A6" if fin else RED),True,False,MONO)]],anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.RIGHT)
dy+=5*(lh+0.06)+0.15
# 2x2 gates (left text + right chart)
gates=[("Gate 2 · selection skill","62th pct","fig_random_entry_mc.png","1,000 random portfolios vs the rule. Indistinguishable from random picking, p = 0.378."),
       ("Mechanism check · expiry timing","+1.16% gap","d5_expiry_premium.png","Signals before expiry return more — a real footprint. Clustered by 84 expiry days: not significant (p=0.10)."),
       ("Gate 3 · regime “edge”","p 0.05→0.47","h2_small_vs_large_sample.png","High vs Low-VIX win rate. On all signals the gap vanishes (z=0.065) — a small-sample artifact."),
       ("Gate 4 · ML drivers","1 / 3","d4_h3_feature_heatmap.png","● = feature in a model’s top-3. The rule features are not the drivers; VIX & MA-distance are.")]
cellw=(DW-0.4)/2; cellh=5.2
for k,(kk,bn,img,cap) in enumerate(gates):
    gx=DX+(k%2)*(cellw+0.4); gyy=dy+(k//2)*(cellh+0.3)
    panel(gx,gyy,cellw,cellh)
    tw=cellw*0.40
    yk=kick(gx+0.22,gyy+0.2,tw-0.3,kk)
    text(gx+0.22,yk,tw-0.3,1.0,[[(bn,38,RED,True,False,SERIF)]])
    text(gx+0.22,yk+1.1,tw-0.3,cellh-1.5,[[(cap,14,MUTED)]])
    pic(img,gx+tw+0.1,gyy+0.4,cellw-tw-0.3,cellh-0.9)

# ================= E. SO WHAT =================
ey=midy+0.85
panel(EX,ey,EW,3.2)
yy=kick(EX+0.25,ey+0.18,EW-0.5,"Interpretation")
text(EX+0.25,yy,EW-0.5,0.9,[[("A real mechanism ≠ a tradeable edge",26,NAVY,True,False,SERIF)]])
text(EX+0.25,yy+0.95,EW-0.5,1.5,[
    [("• The effect is real, but lives in index-option hedging flows — not single-stock RSI.",14,INK)],
    [("• Large-caps are heavily arbitraged; documented edges decay (McLean & Pontiff, 2016).",14,INK)],
    [("• A positive long-only backtest in a bull market is the default, not an achievement.",14,INK)]])
ey+=3.65
panel(EX,ey,EW,7.0)
yy=kick(EX+0.25,ey+0.18,EW-0.5,"The mechanism survives — but not for you")
text(EX+0.25,yy,EW-0.5,1.6,[
    [("▸ Best case, you only match the free ~12%/yr an ETF gives in one click.",14,INK)],
    [("▸ Worst case, concentrated single-name risk (54% drawdown) that underperforms it.",14,INK)],
    [("Rational move: own the baseline cheaply.",14,INK,True)]])
pic("e_fig1_risk_comparison.png",EX+0.25,ey+1.9,EW-0.5,4.9)
ey+=7.45
panel(EX,ey,EW,2.0,fill=PAPER,line=GOLD)
text(EX+0.3,ey+0.12,EW-0.5,1.3,[[("“Don’t confuse a bull market with brains.”",24,NAVY,True,False,SERIF)],[("The rule’s “profit” was the tide, not the trader.",14,MUTED,False,True)]],anchor=MSO_ANCHOR.MIDDLE)
ey+=2.45
panel(EX,ey,EW,3.6)
yy=kick(EX+0.25,ey+0.18,EW-0.5,"Decision flow")
pic("e_fig2_decision_flow.png",EX+0.25,ey+0.6,EW-0.5,2.9)
ey+=4.05
th=panel(EX,ey,EW,5.0,fill=NAVY,line=NAVY)
text(EX+0.25,ey+0.18,EW-0.5,0.35,[[("THE HELP THIS BRINGS — OUR PUBLIC TOOL",13,"7DE0A6",True,False,MONO)]])
text(EX+0.25,ey+0.55,EW-0.5,0.6,[[("“Edge or Beta?”",26,"FFFFFF",True,False,SERIF)]])
text(EX+0.25,ey+1.2,EW-0.5,1.2,[
    [("Your rule  +9.1%    ·    Random stocks  +9.4%",16,"CFE0F5",False,False,MONO)],
    [("Random ETF timing  +10.2%    ·    Hold SPY  +11.9%",16,"CFE0F5",False,False,MONO)],
    [("→ no evidence it beats a low-cost ETF.",16,"7DE0A6",True)]])
text(EX+0.25,ey+2.7,EW-0.5,0.5,[[("niagaradataanalyst.com/edge-or-beta",15,"CFE0F5",False,False,MONO)]])

# ================= F. FOOTER =================
fy=32.5
rect(M,fy,46.4,0.05,fill=NAVY)
boxes=[("FINAL VERDICT",RED,["H1 ✗   H2 ✗   H3 ✗","No edge — it did not beat a cheap S&P 500 index fund.","Don’t trade it as skill; just own the index."]),
       ("WHY IT MATTERS",GREEN,["Rigorous, honest negative result","Separates mechanism from tradeability","Exposes a sample-size illusion (22 → 207)","Turns the null into a public tool"]),
       ("LIMITATIONS & FUTURE RESEARCH",AMBER,["Current-constituent survivorship (not point-in-time)","Single-stock proxy ≠ direct options-level test","Future: point-in-time universe · options-level test · live paper-trading"]),
       ("DATA SOURCES",CYAN,["Daily OHLCV, all S&P 500 — Stooq + Alpaca API","VIX (VIXCLS) — FRED, St. Louis Fed","FOMC — Federal Reserve · Earnings — Nasdaq","Constituents — Wikipedia / SP500 list"]),
       ("REFERENCES",NAVY,["Baltussen, van Dijk & Zhu (2025) · Brock et al. (1992) · Fama (1991) · Golez & Jackwerth (2012) · Kim, Nelson & Startz (1991) · McLean & Pontiff (2016) · Ni et al. (2005) · Poterba & Summers (1988) · Sullivan, Timmermann & White (1999) · Wilder (1978)."])]
fw=[7.6,8.4,9.2,9.2,11.0]; fx=M
for (title,col,items),wdt in zip(boxes,fw):
    text(fx,fy+0.2,wdt,0.4,[[(title,17,col,True,False,MONO)]])
    paras=[[("• "+it if title!="REFERENCES" else it,13 if title!="REFERENCES" else 11,INK if title!="REFERENCES" else MUTED)] for it in items]
    text(fx,fy+0.62,wdt,2.4,paras)
    fx+=wdt+0.4

prs.save(str(OUT))
print("wrote", OUT)
